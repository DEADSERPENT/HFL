"""Edit RMTC_Preseantation.pptx in-place — content update only."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PATH = "/home/deadserpent/HFL/RMTC_Preseantation.pptx"

DARK_BLUE  = RGBColor(0x0D, 0x2B, 0x55)
MID_BLUE   = RGBColor(0x1A, 0x5E, 0x9A)
LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xF7)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE     = RGBColor(0xF4, 0x7F, 0x20)
GREY       = RGBColor(0x33, 0x33, 0x33)
GREEN      = RGBColor(0x1E, 0x8B, 0x4C)
RED        = RGBColor(0xC0, 0x20, 0x20)

prs = Presentation(PATH)
BLANK = prs.slide_layouts[6]


# ── helpers ───────────────────────────────────────────────────────────────────

def clear_slide(slide):
    """Remove every shape from a slide."""
    sp_list = slide.shapes._spTree
    to_del = list(slide.shapes)
    for shape in to_del:
        sp = shape._element
        sp.getparent().remove(sp)

def box(slide, l, t, w, h, fill=DARK_BLUE):
    s = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    return s

def txt(slide, text, l, t, w, h,
        size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def header(slide, title, sub=""):
    box(slide, 0, 0, 13.33, 1.25, DARK_BLUE)
    txt(slide, title, 0.3, 0.08, 12.7, 0.65,
        size=26, bold=True, color=WHITE)
    if sub:
        txt(slide, sub, 0.3, 0.72, 12.7, 0.45,
            size=13, color=LIGHT_BLUE)

def footer(slide):
    box(slide, 0, 7.1, 13.33, 0.4, DARK_BLUE)
    txt(slide, "Samartha H V  |  MIT Bengaluru  |  M.Tech CSE  |  Regd. 251580130019  |  April 24, 2026",
        0.3, 7.12, 12.7, 0.35, size=10,
        color=RGBColor(0xAA,0xCC,0xEE))

def section_bar(slide, label, l, t, w, color=MID_BLUE):
    box(slide, l, t, w, 0.36, color)
    txt(slide, label, l+0.1, t+0.04, w-0.2, 0.28,
        size=12, bold=True, color=WHITE)

def bullets(slide, items, l, t, w, h,
            size=13, color=GREY, spacing=0.42, sym="->"):
    top = t
    for item in items:
        txt(slide, f"  {sym}  {item}", l, top, w, spacing,
            size=size, color=color)
        top += spacing
        if top > t + h:
            break

def kv_row(slide, key, val, l, t, w, h,
           kw=2.8, key_color=MID_BLUE, val_color=GREY, bg=WHITE, size=12):
    box(slide, l, t, w, h, bg)
    txt(slide, key, l+0.1, t+0.04, kw-0.1, h-0.08,
        size=size, bold=True, color=key_color)
    txt(slide, val, l+kw, t+0.04, w-kw-0.1, h-0.08,
        size=size, color=val_color)


# ==============================================================================
# SLIDE 1 — TITLE  (keep existing, just update date text)
# ==============================================================================
# Already looks good — skip


# ==============================================================================
# SLIDE 2 — INTRODUCTION
# ==============================================================================
slide = prs.slides[1]
clear_slide(slide)
box(slide, 0, 0, 13.33, 7.5, RGBColor(0xF5,0xF8,0xFC))
header(slide, "Introduction", "Problem, Motivation & Proposed Solution")

# LEFT — problem
section_bar(slide, "Problem Statement", 0.3, 1.35, 6.1, RED)
items_prob = [
    "2.5 Billion healthcare IoT devices by 2027",
    "Raw ECG & X-ray data is PRIVATE — cannot leave hospital",
    "Centralised ML violates HIPAA / GDPR patient privacy",
    "High communication cost uploading raw data to cloud",
    "Latency unacceptable for real-time clinical decisions",
]
bullets(slide, items_prob, 0.35, 1.78, 5.95, 2.5, size=13)

# RIGHT — solution
section_bar(slide, "Proposed Solution — Hierarchical Federated Learning",
            6.6, 1.35, 6.4, GREEN)
items_sol = [
    "Devices train locally — raw data NEVER shared",
    "Only model weights (gradients) are transmitted",
    "Two-tier hierarchy:  Device -> Edge -> Cloud",
    "DP-SGD adds mathematical privacy guarantee (epsilon=8.0)",
    "Top-k sparsification reduces communication by 76.4%",
]
bullets(slide, items_sol, 6.65, 1.78, 6.25, 2.5, size=13)

# Domain badges
section_bar(slide, "Research Domain — Healthcare Center (Primary Focus)",
            0.3, 4.45, 12.7, DARK_BLUE)
for i, (lbl, sub, col) in enumerate([
    ("Modality A\nPTB-XL ECG", "12-lead Signal\n5,000 records", MID_BLUE),
    ("Modality B\nKermany CXR", "Chest X-Ray Image\n3,000 images", GREEN),
    ("Modality C\nTabular", "PTB-XL Metadata\n5 features", ORANGE),
    ("FL Setup\n6 Devices", "2 Edge Servers\nDirichlet alpha=0.5", DARK_BLUE),
]):
    x = 0.3 + i * 3.25
    box(slide, x, 4.88, 3.1, 1.5, col)
    txt(slide, lbl, x+0.1, 4.92, 2.9, 0.65,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, sub, x+0.1, 5.57, 2.9, 0.7,
        size=12, color=WHITE, align=PP_ALIGN.CENTER)

footer(slide)


# ==============================================================================
# SLIDE 3 — LITERATURE SURVEY & GAPS
# ==============================================================================
slide = prs.slides[2]
clear_slide(slide)
box(slide, 0, 0, 13.33, 7.5, RGBColor(0xF5,0xF8,0xFC))
header(slide, "Literature Survey & Research Gaps", "What exists — What is missing")

# Table header
section_bar(slide, "Existing Work Review", 0.3, 1.35, 12.7, DARK_BLUE)
cols = [("Reference", 3.2), ("Method", 4.5), ("Limitation", 4.7)]
cx = 0.3
for h, w in cols:
    box(slide, cx, 1.76, w, 0.38, MID_BLUE)
    txt(slide, h, cx+0.08, 1.78, w-0.1, 0.34,
        size=12, bold=True, color=WHITE)
    cx += w + 0.05

rows = [
    ("McMahan et al. 2017 (FedAvg)",
     "Single-tier federated averaging",
     "No hierarchy, no privacy, single modality only"),
    ("Li et al. 2020 (FedProx)",
     "Proximal term for non-IID heterogeneity",
     "No multimodal support, no DP guarantee"),
    ("Abadi et al. 2016 (DP-SGD)",
     "Differential privacy for deep learning",
     "Not applied in federated multimodal healthcare"),
    ("Qayyum et al. 2022 (FL Survey)",
     "Healthcare FL survey — medical imaging",
     "Single modality only, ECG+X-ray fusion missing"),
    ("MOON — Li et al. 2021",
     "Contrastive FL for non-IID clients",
     "No DP guarantee, flat architecture only"),
]
for ri, row in enumerate(rows):
    y = 2.2 + ri * 0.46
    bg = WHITE if ri % 2 == 0 else LIGHT_BLUE
    cx = 0.3
    for cell, (_, w) in zip(row, cols):
        box(slide, cx, y, w, 0.44, bg)
        txt(slide, cell, cx+0.08, y+0.05, w-0.12, 0.36,
            size=11, color=GREY)
        cx += w + 0.05

# Gaps + Objectives
section_bar(slide, "Research Gaps Identified", 0.3, 4.6, 6.2, ORANGE)
bullets(slide, [
    "No FL system fuses ECG + Chest X-ray simultaneously",
    "DP-SGD + Top-k compression not combined in hierarchical FL",
    "Non-IID multimodal healthcare partitioning not studied",
], 0.3, 5.0, 6.2, 1.8, size=13)

section_bar(slide, "Research Objectives", 6.8, 4.6, 6.2, GREEN)
bullets(slide, [
    "Design two-tier HFL for multimodal healthcare IoT",
    "Guarantee patient privacy via DP-SGD (epsilon=8.0)",
    "Reduce communication by >70% without accuracy loss",
], 6.8, 5.0, 6.2, 1.8, size=13)

footer(slide)


# ==============================================================================
# SLIDE 4 — ARCHITECTURE
# ==============================================================================
slide = prs.slides[3]
clear_slide(slide)
box(slide, 0, 0, 13.33, 7.5, RGBColor(0xF5,0xF8,0xFC))
header(slide, "System Architecture",
       "Two-Tier Hierarchical Federated Learning — Healthcare Domain")

# 3 tier boxes
for x, w, col, title, lines in [
    (0.3,  3.9, DARK_BLUE, "TIER 1 — IoT Devices",
     ["6 wearable / bedside sensors",
      "Local DP-SGD training",
      "Top-k gradient compression",
      "ECG + X-ray stays on device"]),
    (4.45, 4.4, MID_BLUE,  "TIER 2 — Edge Servers",
     ["2 ward-level gateways",
      "Weighted FedAvg aggregation",
      "FedConform-HC calibration",
      "Low-latency local aggregation"]),
    (9.1,  3.9, GREEN,     "TIER 3 — Cloud Server",
     ["Global model aggregation",
      "Privacy budget tracking (eps,delta)",
      "Broadcast updated model to edges",
      "Checkpoint best model"]),
]:
    box(slide, x, 1.35, w, 2.3, col)
    txt(slide, title, x+0.12, 1.38, w-0.2, 0.4,
        size=13, bold=True, color=WHITE)
    for i, line in enumerate(lines):
        txt(slide, f"->  {line}", x+0.12, 1.82+i*0.44, w-0.2, 0.42,
            size=12, color=WHITE)

txt(slide, "->", 4.18, 2.1, 0.4, 0.5, size=24, bold=True,
    color=ORANGE, align=PP_ALIGN.CENTER)
txt(slide, "->", 8.87, 2.1, 0.4, 0.5, size=24, bold=True,
    color=ORANGE, align=PP_ALIGN.CENTER)

# Model 1
section_bar(slide, "Model 1 — HFL-MM-HC  (Baseline)", 0.3, 3.85, 6.1, MID_BLUE)
kv_data_1 = [
    ("ECG Encoder",   "1D-CNN + Bidirectional GRU  ->  [B, 256]"),
    ("Image Encoder", "MobileNetV3-Small (pretrained)  ->  [B, 576]"),
    ("Fusion",        "Late fusion FC head + GroupNorm"),
    ("Privacy",       "DP-SGD  |  epsilon=8.0  |  delta=1e-5"),
    ("Compression",   "Top-k 80% sparse + 8-bit quantisation"),
    ("Parameters",    "2.15M total  |  2.0M trainable"),
]
for i, (k, v) in enumerate(kv_data_1):
    y = 4.27 + i * 0.45
    bg = WHITE if i % 2 == 0 else LIGHT_BLUE
    kv_row(slide, k, v, 0.3, y, 6.1, 0.43,
           kw=2.2, bg=bg, size=11)

# Model 2
section_bar(slide, "Model 2 — PHANTOM-FL  (5 Innovations)", 6.8, 3.85, 6.2, ORANGE)
kv_data_2 = [
    ("CMGA",    "Cross-Modal Gated Attention (ECG <-> X-ray)  [IP-3]"),
    ("SADP",    "Self-Adaptive Differential Privacy           [IP-4]"),
    ("ASWA",    "Adaptive Server-side Weight Aggregation      [IP-5]"),
    ("FedKD-E", "Federated Knowledge Distillation at Edge     [IP-6]"),
    ("pFAL",    "Personalised Federated Adapter Learning      [IP-7]"),
    ("IPs",     "10 IP claims filed  |  7 novelty contributions"),
]
for i, (k, v) in enumerate(kv_data_2):
    y = 4.27 + i * 0.45
    bg = WHITE if i % 2 == 0 else LIGHT_BLUE
    kv_row(slide, k, v, 6.8, y, 6.2, 0.43,
           kw=1.8, bg=bg, size=11)

footer(slide)


# ==============================================================================
# SLIDE 5 — DATASETS
# ==============================================================================
slide = prs.slides[4]
clear_slide(slide)
box(slide, 0, 0, 13.33, 7.5, RGBColor(0xF5,0xF8,0xFC))
header(slide, "Datasets Used",
       "Real Clinical Datasets — Peer-Reviewed & Publicly Available")

# PTB-XL
section_bar(slide, "Modality A — PTB-XL ECG  (Signal / Time-series)", 0.3, 1.35, 6.1, MID_BLUE)
ptbxl = [
    ("Source",      "PhysioNet  /  Nature Scientific Data 2020"),
    ("Full size",   "21,799 ECG records — 12 leads at 100 Hz"),
    ("Used",        "5,000 records — stratified 1,000 per class"),
    ("Classes",     "NORM  |  MI  |  STTC  |  CD  |  HYP"),
    ("Tensor",      "[5000, 12, 1000]  float32"),
    ("Tabular",     "Age, sex, height, weight, strat_fold (5 features)"),
    ("Why",         "Standard benchmark for cardiac FL research"),
]
for i, (k, v) in enumerate(ptbxl):
    y = 1.78 + i * 0.46
    bg = WHITE if i % 2 == 0 else LIGHT_BLUE
    kv_row(slide, k, v, 0.3, y, 6.1, 0.44, kw=1.8, bg=bg, size=11)

# CXR
section_bar(slide, "Modality B — Kermany Chest X-Ray  (Image)", 6.8, 1.35, 6.2, GREEN)
cxr = [
    ("Source",   "Cell 2018 — Kermany et al.  (top-tier citation)"),
    ("Full size","5,856 JPEG images"),
    ("Used",     "3,000 images — balanced 1,500 per class"),
    ("Classes",  "NORMAL (0)  |  PNEUMONIA (1)"),
    ("Format",   "JPEG -> RGB -> [3, 224, 224]  float32"),
    ("Tensor",   "[3000, 3, 224, 224]  float32"),
    ("Why",      "Complements ECG in realistic ward scenario"),
]
for i, (k, v) in enumerate(cxr):
    y = 1.78 + i * 0.46
    bg = WHITE if i % 2 == 0 else LIGHT_BLUE
    kv_row(slide, k, v, 6.8, y, 6.2, 0.44, kw=1.8, bg=bg, size=11)

# Partition summary
section_bar(slide, "FL Non-IID Partition  (Dirichlet alpha=0.5)  |  Avg TVD=0.4556",
            0.3, 5.15, 12.7, DARK_BLUE)
hdrs = ["Device","Edge","Train","Val","Test","Class Distribution"]
wx   = [1.5, 1.2, 1.2, 1.0, 1.0, 6.6]
rows_p = [
    ("Device 00","Edge 0","557","119","120","Heavy CD  (595 samples)"),
    ("Device 01","Edge 0","221","47","49","Mixed NORM+STTC"),
    ("Device 02","Edge 0","694","148","150","Heavy MI+STTC"),
    ("Device 03","Edge 1","763","163","164","Heavy HYP  (813 samples)"),
    ("Device 04","Edge 1","536","114","116","Heavy NORM+MI"),
    ("Device 05","Edge 1","727","155","157","Mixed MI+CD+NORM"),
]
cx_start = 0.3
for hi, (h, w) in enumerate(zip(hdrs, wx)):
    box(slide, cx_start, 5.56, w, 0.36, MID_BLUE)
    txt(slide, h, cx_start+0.04, 5.58, w-0.06, 0.32,
        size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    cx_start += w + 0.02

for ri, row in enumerate(rows_p):
    y = 5.94 + ri * 0.34
    bg = WHITE if ri % 2 == 0 else LIGHT_BLUE
    cx_start = 0.3
    for ci, (cell, w) in enumerate(zip(row, wx)):
        box(slide, cx_start, y, w, 0.32, bg)
        col = GREEN if cell in ("Edge 0","Edge 1") else GREY
        txt(slide, cell, cx_start+0.04, y+0.04, w-0.06, 0.26,
            size=9, color=col, align=PP_ALIGN.CENTER)
        cx_start += w + 0.02

footer(slide)


# ==============================================================================
# SLIDE 6 — METHODOLOGY
# ==============================================================================
slide = prs.slides[5]
clear_slide(slide)
box(slide, 0, 0, 13.33, 7.5, RGBColor(0xF5,0xF8,0xFC))
header(slide, "Methodology & Algorithm",
       "Two-Tier FedAvg  +  DP-SGD  +  Top-k Gradient Compression")

section_bar(slide, "Training Algorithm — HFL-MM-HC  (20 Global Rounds x 5 Edge Rounds)",
            0.3, 1.35, 8.2, DARK_BLUE)
steps = [
    ("1", "Initialise",   "Cloud broadcasts global model to 2 edge servers"),
    ("2", "Distribute",   "Each edge server sends model to its 3 devices"),
    ("3", "Local Train",  "Device runs DP-SGD with Gaussian noise (epsilon=8.0, delta=1e-5)"),
    ("4", "Compress",     "Top-k sparsification (80% sparse) + 8-bit quantisation"),
    ("5", "Edge FedAvg",  "Edge server: weighted FedAvg of compressed gradients"),
    ("6", "Conformal",    "FedConform-HC calibrates uncertainty prediction sets"),
    ("7", "Cloud Agg",    "Cloud aggregates 2 edge models -> new global model"),
    ("8", "Evaluate",     "Measure AUC, Accuracy, epsilon-spent on validation set"),
    ("9", "Repeat",       "Steps 2-8 for 5 edge rounds per global round -> 20 rounds"),
]
for i, (num, name, desc) in enumerate(steps):
    y = 1.78 + i * 0.52
    box(slide, 0.3,  y, 0.48, 0.48, ORANGE)
    txt(slide, num, 0.3, y+0.06, 0.48, 0.36,
        size=14, bold=True, align=PP_ALIGN.CENTER)
    box(slide, 0.82, y, 1.7,  0.48, MID_BLUE)
    txt(slide, name, 0.88, y+0.07, 1.6, 0.36, size=11, bold=True)
    bg = WHITE if i % 2 == 0 else LIGHT_BLUE
    box(slide, 2.56, y, 6.0,  0.48, bg)
    txt(slide, desc, 2.65, y+0.07, 5.85, 0.36, size=11, color=GREY)

# Right panel — DP params
section_bar(slide, "DP-SGD Parameters", 8.85, 1.35, 4.15, DARK_BLUE)
dp = [
    ("Privacy budget (epsilon)", "8.0"),
    ("Delta (delta)",            "1 x 10^-5"),
    ("Noise sigma (auto)",       "0.5655"),
    ("Gradient clip C",          "1.0"),
    ("Mechanism",                "Gaussian — Opacus"),
    ("Guarantee",                "Mathematical"),
]
for i, (k, v) in enumerate(dp):
    y = 1.78 + i * 0.52
    bg = WHITE if i % 2 == 0 else LIGHT_BLUE
    box(slide, 8.85, y, 4.15, 0.48, bg)
    txt(slide, k, 8.95, y+0.07, 2.5, 0.36, size=11, bold=True, color=MID_BLUE)
    txt(slide, v, 11.5, y+0.07, 1.4, 0.36, size=11, bold=True, color=ORANGE,
        align=PP_ALIGN.CENTER)

section_bar(slide, "Compression", 8.85, 6.52, 4.15, MID_BLUE)
comp = [
    "Top-k sparsification  ->  80% sparse",
    "8-bit quantisation    ->  4x size reduction",
    "Error feedback buffer ->  no accuracy loss",
    "Total comm. saving    ->  76.4%  (confirmed)",
]
for i, c in enumerate(comp):
    txt(slide, f"->  {c}", 8.95, 6.92+i*0.0, 4.0, 0.36, size=11, color=GREY)
bullets(slide, comp, 8.95, 6.92, 4.0, 1.5, size=11, spacing=0.35)

footer(slide)


# ==============================================================================
# SLIDE 7 — RESULTS & EXPECTED OUTCOMES
# ==============================================================================
slide = prs.slides[6]
clear_slide(slide)
box(slide, 0, 0, 13.33, 7.5, RGBColor(0xF5,0xF8,0xFC))
header(slide, "Results & Expected Outcomes",
       "Simulation Confirmed  |  Model Training In Progress")

# Status banner
box(slide, 0.3, 1.35, 12.7, 0.42, ORANGE)
txt(slide,
    "SIMULATION RESULTS: Confirmed via NS-3 v3.46.1 + CloudSim Plus v8.5.6    "
    "MODEL RESULTS: Pipeline complete — training in progress",
    0.5, 1.38, 12.4, 0.36, size=11, bold=True, color=WHITE)

# Confirmed simulation results
section_bar(slide, "Simulation Results  (CONFIRMED)", 0.3, 1.87, 6.1, DARK_BLUE)
sim = [
    ("Communication Reduction", "76.4%",  "> 70%",  True),
    ("Energy Saving",           "74.9%",  "> 70%",  True),
    ("Network Reliability",     "99.05%", "> 99%",  True),
    ("Round-trip Latency",      "< 120ms","< 150ms", True),
]
for ci, (h, w) in enumerate(zip(
    ["Metric", "Achieved", "Target"],
    [3.1, 1.5, 1.3])):
    x = 0.3 + sum([3.1,1.5,1.3][:ci]) + ci*0.07
    box(slide, x, 2.3, w, 0.36, MID_BLUE)
    txt(slide, h, x+0.05, 2.32, w-0.1, 0.32,
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri, (m, a, t, ok) in enumerate(sim):
    y = 2.68 + ri * 0.46
    bg = WHITE if ri % 2 == 0 else LIGHT_BLUE
    for ci, (cell, w) in enumerate(zip([m,a,t],
                                        [3.1,1.5,1.3])):
        x = 0.3 + sum([3.1,1.5,1.3][:ci]) + ci*0.07
        box(slide, x, y, w, 0.44, bg)
        color = GREEN if (ci == 1 and ok) else GREY
        txt(slide, cell, x+0.07, y+0.07, w-0.1, 0.32,
            size=12, bold=(ci==1), color=color,
            align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    # Exceeded badge
    box(slide, 5.6, y, 0.78, 0.44, GREEN)
    txt(slide, "MET", 5.6, y+0.07, 0.78, 0.32,
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Expected model results
section_bar(slide, "Expected Model Results  (Training In Progress)", 6.8, 1.87, 6.2, DARK_BLUE)
exp = [
    ("ECG 5-class Accuracy",  "82 - 88%",
     "Based on PTB-XL FL literature"),
    ("Macro AUC (5-class)",   "0.90 - 0.94",
     "Comparable multimodal FL works"),
    ("Weighted F1-Score",     "0.83 - 0.89",
     "Non-IID federated benchmark"),
    ("FL Convergence",        "~15 rounds",
     "Standard FedAvg convergence"),
    ("Final Privacy budget",  "epsilon <= 8.0",
     "DP-SGD configured and running"),
]
for ci, (h, w) in enumerate(zip(
    ["Metric","Expected","Basis"],
    [2.5, 1.5, 2.0])):
    x = 6.8 + sum([2.5,1.5,2.0][:ci]) + ci*0.07
    box(slide, x, 2.3, w, 0.36, MID_BLUE)
    txt(slide, h, x+0.05, 2.32, w-0.1, 0.32,
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri, (m, v, basis) in enumerate(exp):
    y = 2.68 + ri * 0.46
    bg = WHITE if ri % 2 == 0 else LIGHT_BLUE
    for ci, (cell, w) in enumerate(zip([m, v, basis],
                                        [2.5, 1.5, 2.0])):
        x = 6.8 + sum([2.5,1.5,2.0][:ci]) + ci*0.07
        box(slide, x, y, w, 0.44, bg)
        color = ORANGE if ci == 1 else GREY
        txt(slide, cell, x+0.07, y+0.07, w-0.1, 0.32,
            size=11, bold=(ci==1), color=color,
            align=PP_ALIGN.CENTER if ci == 1 else PP_ALIGN.LEFT)

# Bottom note
box(slide, 0.3, 5.1, 12.7, 0.36, LIGHT_BLUE)
txt(slide,
    "Expected values are grounded in:  Wagner et al. 2020 (PTB-XL FL)  |  "
    "Nguyen et al. 2022 (Healthcare FL)  |  McMahan et al. 2017 (FedAvg baseline)",
    0.5, 5.12, 12.5, 0.32, size=10, color=DARK_BLUE)

# Implementation status strip
section_bar(slide, "Implementation Status", 0.3, 5.55, 12.7, DARK_BLUE)
statuses = [
    ("HFL Architecture",    "DONE", GREEN),
    ("Dataset Pipeline",    "DONE", GREEN),
    ("Model Code (GPU)",    "DONE", GREEN),
    ("NS-3 Simulation",     "DONE", GREEN),
    ("CloudSim Simulation", "DONE", GREEN),
    ("Model Training",      "IN PROGRESS", ORANGE),
    ("Final Paper",         "DRAFTING", ORANGE),
]
for i, (comp, stat, col) in enumerate(statuses):
    x = 0.3 + i * 1.86
    box(slide, x, 5.97, 1.78, 0.78, WHITE)
    txt(slide, comp, x+0.06, 5.99, 1.66, 0.38, size=9, bold=True, color=DARK_BLUE)
    box(slide, x+0.08, 6.4, 1.62, 0.3, col)
    txt(slide, stat, x+0.08, 6.41, 1.62, 0.28, size=8, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

footer(slide)


# ==============================================================================
# SLIDE 8 — IP & NOVELTY
# ==============================================================================
slide = prs.slides[7]
clear_slide(slide)
box(slide, 0, 0, 13.33, 7.5, RGBColor(0xF5,0xF8,0xFC))
header(slide, "Novelty & IP Contributions",
       "10 Intellectual Property Claims  +  7 Novelty Contributions")

section_bar(slide, "IP Claims — Healthcare Domain", 0.3, 1.35, 12.7, DARK_BLUE)
ips = [
    ("IP-1",  "FedMamba-HC",    "State-space ECG encoder for federated healthcare FL"),
    ("IP-2",  "FedMedLoRA",     "Low-rank adaptation for medical FL fine-tuning"),
    ("IP-3",  "ClinicalCMGA",   "Cross-modal gated attention — ECG + X-ray fusion"),
    ("IP-4",  "SADP",           "Self-adaptive differential privacy noise scheduler"),
    ("IP-5",  "ASWA",           "Adaptive server-side weight aggregation algorithm"),
    ("IP-6",  "FedKD-E",        "Edge-level federated knowledge distillation"),
    ("IP-7",  "pFAL",           "Personalised federated adapter learning"),
    ("IP-8",  "HFL-MM",         "Multimodal hierarchical FL baseline framework"),
    ("IP-9",  "FedConform-HC",  "Federated conformal prediction for healthcare"),
    ("IP-10", "PHANTOM-FL",     "Full 5-innovation privacy-aware multimodal FL system"),
]
for ci, (h, w) in enumerate(zip(["IP","Name","Contribution"],[0.8,2.2,9.5])):
    x = 0.3 + sum([0.8,2.2,9.5][:ci]) + ci*0.05
    box(slide, x, 1.78, w, 0.36, MID_BLUE)
    txt(slide, h, x+0.06, 1.8, w-0.1, 0.32,
        size=12, bold=True, color=WHITE)

for ri, (ip, name, contrib) in enumerate(ips):
    y = 2.16 + ri * 0.45
    bg = WHITE if ri % 2 == 0 else LIGHT_BLUE
    col = ORANGE if ri < 5 else MID_BLUE
    for ci, (cell, w) in enumerate(zip([ip,name,contrib],[0.8,2.2,9.5])):
        x = 0.3 + sum([0.8,2.2,9.5][:ci]) + ci*0.05
        box(slide, x, y, w, 0.43, col if ci==0 else bg)
        txt(slide, cell, x+0.06, y+0.06, w-0.1, 0.33,
            size=11,
            bold=(ci==0),
            color=WHITE if ci==0 else (DARK_BLUE if ci==1 else GREY),
            align=PP_ALIGN.CENTER if ci==0 else PP_ALIGN.LEFT)

# Target journal
box(slide, 0.3, 6.72, 12.7, 0.46, DARK_BLUE)
txt(slide,
    "Target Publication:  IEEE Transactions on Neural Networks and Learning Systems (TNNLS)  |  "
    "Q1 Journal  |  Impact Factor > 10",
    0.5, 6.76, 12.5, 0.38, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

footer(slide)


# ==============================================================================
# SLIDE 9 — CONCLUSION & REFERENCES
# ==============================================================================
slide = prs.slides[8]
clear_slide(slide)
box(slide, 0, 0, 13.33, 7.5, RGBColor(0xF5,0xF8,0xFC))
header(slide, "Conclusion & References", "")

# Abstract
section_bar(slide, "Abstract", 0.3, 1.35, 12.7, MID_BLUE)
box(slide, 0.3, 1.75, 12.7, 1.0, WHITE)
txt(slide,
    "This research proposes a Hierarchical Federated Learning (HFL) framework for privacy-aware "
    "multimodal IoT systems in healthcare. A 12-lead ECG encoder (1D-CNN+GRU) is combined with "
    "a chest X-ray encoder (MobileNetV3-Small) under a two-tier federated architecture across "
    "6 IoT devices and 2 edge servers. DP-SGD (epsilon=8.0) provides mathematical privacy "
    "guarantees while Top-k sparsification achieves 76.4% communication reduction. Simulation "
    "results confirm 74.9% energy saving and 99.05% network reliability, both exceeding design "
    "targets. Ten IP claims and seven novelty contributions are documented.",
    0.45, 1.8, 12.4, 0.92, size=12, color=GREY)

# Conclusion + Future Work
section_bar(slide, "Conclusion", 0.3, 2.88, 6.1, GREEN)
bullets(slide, [
    "Complete two-tier HFL framework for healthcare IoT",
    "Multimodal fusion: ECG + Chest X-ray under FL",
    "Mathematical privacy guarantee via DP-SGD",
    "76.4% communication reduction  (target exceeded)",
    "99.05% network reliability  (target exceeded)",
    "10 IP claims + 7 novelty contributions filed",
], 0.3, 3.28, 6.1, 2.7, size=12, spacing=0.44)

section_bar(slide, "Future Work", 6.8, 2.88, 6.2, ORANGE)
bullets(slide, [
    "Complete model training — collect AUC, F1 scores",
    "Implement PHANTOM-FL full 5-innovation pipeline",
    "Extend to 3rd modality (clinical report text)",
    "Deploy on Raspberry Pi edge hardware prototype",
    "Submit to IEEE TNNLS / TPDS journal",
], 6.8, 3.28, 6.2, 2.7, size=12, spacing=0.44)

# References
section_bar(slide, "References", 0.3, 6.08, 12.7, DARK_BLUE)
refs = [
    "[1] McMahan et al. — Communication-Efficient Learning of Deep Networks from Decentralized Data. AISTATS 2017.",
    "[2] Abadi et al.  — Deep Learning with Differential Privacy. ACM CCS 2016.",
    "[3] Wagner et al. — PTB-XL, A Large Publicly Available ECG Dataset. Nature Scientific Data 2020.",
    "[4] Kermany et al.— Identifying Medical Diagnoses & Treatable Diseases using Deep-Learning. Cell 2018.",
]
for i, r in enumerate(refs):
    txt(slide, r, 0.4, 6.5+i*0.23, 12.7, 0.22, size=9,
        italic=True, color=GREY)

footer(slide)


# ── SAVE ──────────────────────────────────────────────────────────────────────
prs.save(PATH)
print(f"Saved in-place: {PATH}")
print(f"Total slides  : {len(prs.slides)}")
